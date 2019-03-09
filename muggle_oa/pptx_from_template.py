# -*- coding: utf-8 -*-
# author = 'K0ctr'

from pptx import Presentation
import xlrd


def change_text(powerpoint, old_text, new_text):
    for slide in powerpoint.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run_text = run.text.replace(old_text, new_text)
                    run.text = run_text


def change_text_from_xls():
    xls = xlrd.open_workbook("D:/change_ppt.xlsx")
    sheet = xls.sheet_by_index(0)
    for row in range(1, sheet.nrows):
        ppt = Presentation("D:/mode.pptx")
        for col in range(0, sheet.ncols):
            change_text(ppt, str(sheet.cell_value(0, col)), str(sheet.cell_value(row, col)))

        ppt.save("D:/change_mode.pptx")


change_text_from_xls()
