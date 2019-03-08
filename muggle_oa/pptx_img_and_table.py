# -*- coding: utf-8 -*-
# author = 'K0ctr'

from pptx import Presentation
from pptx.util import Inches

ppt = Presentation()
slide = ppt.slides.add_slide(ppt.slide_layouts[1])

left = top = Inches(0)
width, height = Inches(2.11), Inches(0.33)
pic = slide.shapes.add_picture('D:/banner.jpg', left, top, width, height)

rows = 3
cols = 3
left = top = Inches(1)
width = height = Inches(6)
table = slide.shapes.add_table(rows, cols, left, top, width, height).table
table.columns[0].width = table.columns[1].width = table.columns[2].width = Inches(2)
for row in range(0, rows):
    for col in range(0, cols):
        table.cell(row, col).text = str(row * col)

table.cell(0, 0).merge(table.cell(0, 2))

ppt.save("D:/ppt_img_and_table.pptx")
