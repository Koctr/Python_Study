# -*- coding: utf-8 -*-
# author = 'K0ctr'


from pptx import Presentation
from pptx.util import Inches

ppt = Presentation()
slide = ppt.slides.add_slide(ppt.slide_layouts[1])

left = top = width = height = Inches(1)
textbox1 = slide.shapes.add_textbox(left, top, width, height)
textbox1.text = "first textbox"
left = top = width = height = Inches(2)
textbox2 = slide.shapes.add_textbox(left, top, width, height)
textbox2.text = "second textbox"

ppt.save("D:/ppt_textbox.pptx")
